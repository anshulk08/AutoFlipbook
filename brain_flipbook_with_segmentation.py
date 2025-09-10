import os
import glob
import re
import numpy as np
import nibabel as nib
import matplotlib.pyplot as plt
from matplotlib.gridspec import GridSpec
from datetime import datetime
import warnings    
warnings.filterwarnings('ignore')

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PPTX generation will be disabled.")

class BrainFlipbookGenerator:
    def __init__(self, registered_folder, segmentation_folder=None, output_folder="flipbooks"):
        """
        Initialize the flipbook generator for folder-based timepoint organization
        
        Parameters:
        - registered_folder: Path to folder containing timepoint subfolders (e.g., 5885, 5972, 6070)
                           Also looks for transformed segmentations in <timepoint>/segmentations/ subfolders
        - segmentation_folder: Path to folder containing original segmentation maps (optional fallback)
        - output_folder: Path to output folder for flipbooks
        
        Note: Prioritizes transformed segmentations in registered_folder over original segmentations
        """
        self.registered_folder = registered_folder
        self.segmentation_folder = segmentation_folder
        self.output_folder = output_folder
        self.timepoint_folders = []
        self.contrast_types = []
        
        # Create output folder if it doesn't exist
        os.makedirs(output_folder, exist_ok=True)
        
    def find_timepoint_folders(self):
        """Find all timepoint folders (numeric or date format)"""
        folders = []
        date_pattern = re.compile(r'^(\d{4})-(\d{2})-(\d{2})$')
        
        for item in os.listdir(self.registered_folder):
            item_path = os.path.join(self.registered_folder, item)
            
            if os.path.isdir(item_path):
                # Check if it's a numeric folder name (legacy support)
                if item.isdigit():
                    folders.append({
                        'folder': item_path,
                        'timepoint': int(item),
                        'folder_name': item,
                        'is_date': False
                    })
                # Check if it's a date format (YYYY-MM-DD)
                elif date_pattern.match(item):
                    try:
                        # Parse date and convert to timestamp for sorting
                        date_obj = datetime.strptime(item, '%Y-%m-%d')
                        folders.append({
                            'folder': item_path,
                            'timepoint': int(date_obj.timestamp()),
                            'folder_name': item,
                            'is_date': True,
                            'date_obj': date_obj
                        })
                    except ValueError:
                        continue
        
        # Sort by timepoint (timestamp for dates, numeric for legacy)
        folders.sort(key=lambda x: x['timepoint'])
        self.timepoint_folders = folders
        
        print(f"Found {len(folders)} timepoint folders:")
        for tp in folders:
            print(f"  {tp['folder_name']}: {tp['folder']}")
        
        return folders
    
    def find_contrast_types(self):
        """Find available contrast types across all timepoints"""
        contrast_types = set()
        
        for tp_folder in self.timepoint_folders:
            nifti_files = glob.glob(os.path.join(tp_folder['folder'], "*.nii*"))
            
            for file in nifti_files:
                basename = os.path.basename(file)
                # Parse contrast type from filename patterns:
                # Legacy: "6070_T1_to_5885T1CE.nii"
                # Date-based: "2000-03-25_FL_to_2000-03-25FL.nii.gz"
                match = re.search(r'^(?:\d+|\d{4}-\d{2}-\d{2})_([^_]+)_to_', basename)
                if match:
                    contrast_types.add(match.group(1))
        
        self.contrast_types = sorted(list(contrast_types))
        print(f"Available contrast types: {self.contrast_types}")
        return self.contrast_types
    
    def get_file_for_timepoint_and_contrast(self, timepoint_folder, contrast_type):
        """Get the NIFTI file for a specific timepoint and contrast type"""
        folder_name = os.path.basename(timepoint_folder['folder'])
        
        # Look for file with pattern: {timepoint}_{contrast}_to_{reference}T1CE.nii
        pattern = f"{folder_name}_{contrast_type}_to_*.nii*"
        files = glob.glob(os.path.join(timepoint_folder['folder'], pattern))
        
        if files:
            return files[0]  # Return first match
        else:
            print(f"Warning: No {contrast_type} file found for timepoint {folder_name}")
            return None
    
    def get_segmentation_file(self, timepoint_folder):
        """Get the segmentation file for a specific timepoint
        
        First checks for transformed segmentations in the registered folder,
        then falls back to original segmentations in the segmentation folder.
        """
        if not self.segmentation_folder:
            return None
            
        folder_name = os.path.basename(timepoint_folder['folder'])
        
        # PRIORITY 1: Look for transformed segmentations in the registered folder
        # These are automatically aligned with the registered images
        registered_seg_path = os.path.join(timepoint_folder['folder'], 'segmentations')
        if os.path.exists(registered_seg_path):
            # Look for transformed segmentation file
            transformed_patterns = [
                f"{folder_name}_seg_registered.nii*",
                "*_seg_registered.nii*",
                "*.nii*"
            ]
            
            for pattern in transformed_patterns:
                files = glob.glob(os.path.join(registered_seg_path, pattern))
                if files:
                    print(f"Found transformed segmentation: {files[0]}")
                    return files[0]
        
        # PRIORITY 2: Fall back to original segmentation folder (legacy support)
        # First, try looking for files in a subfolder structure
        seg_folder_path = os.path.join(self.segmentation_folder, folder_name)
        
        if os.path.exists(seg_folder_path):
            # Look in subfolder
            patterns = [
                "*.nii*",  # Any NIFTI file
                "*seg*.nii*",  # Files with 'seg' in name
                "*mask*.nii*",  # Files with 'mask' in name
                "*tumor*.nii*",  # Files with 'tumor' in name
                f"{folder_name}*.nii*"  # Files starting with timepoint number
            ]
            
            for pattern in patterns:
                files = glob.glob(os.path.join(seg_folder_path, pattern))
                if files:
                    print(f"Found segmentation file: {files[0]}")
                    return files[0]
        
        # If no subfolder, look directly in the segmentation folder for files named with timepoint
        direct_patterns = [
            f"{folder_name}_seg.nii*",  # timepoint_seg.nii
            f"{folder_name}_mask.nii*",  # timepoint_mask.nii
            f"{folder_name}_tumor.nii*",  # timepoint_tumor.nii
            f"{folder_name}.nii*",  # timepoint.nii
            f"*{folder_name}*seg*.nii*",  # any file with timepoint and seg
            f"*{folder_name}*.nii*"  # any file with timepoint
        ]
        
        for pattern in direct_patterns:
            files = glob.glob(os.path.join(self.segmentation_folder, pattern))
            if files:
                print(f"Found original segmentation: {files[0]} (note: not transformed)")
                return files[0]
        
        print(f"Warning: No segmentation file found for timepoint {folder_name}")
        return None
    
    def calculate_tumor_volume(self, segmentation_file):
        """Calculate tumor volume from segmentation file"""
        if not segmentation_file or not os.path.exists(segmentation_file):
            return None
        
        try:
            # Load segmentation
            seg_img = nib.load(segmentation_file)
            seg_data = seg_img.get_fdata()
            
            # Handle 4D data
            if len(seg_data.shape) == 4:
                seg_data = seg_data[:, :, :, 0]
            
            # Count tumor voxels (using higher threshold to avoid noise)
            tumor_threshold = 0.5
            tumor_voxels = np.sum(seg_data > tumor_threshold)
            
            # Get voxel dimensions from header
            header = seg_img.header
            voxel_dims = header.get_zooms()[:3]  # x, y, z dimensions in mm
            
            # Calculate volume in mm³
            voxel_volume_mm3 = np.prod(voxel_dims)
            tumor_volume_mm3 = tumor_voxels * voxel_volume_mm3
            
            # Convert to cm³ (mL)
            tumor_volume_ml = tumor_volume_mm3 / 1000.0
            
            return {
                'volume_mm3': tumor_volume_mm3,
                'volume_ml': tumor_volume_ml,
                'voxel_count': tumor_voxels,
                'voxel_dims': voxel_dims
            }
            
        except Exception as e:
            print(f"Error calculating tumor volume: {e}")
            return None
    
    def create_tumor_overlay_colormap(self, base_colormap='gray', tumor_color='red', tumor_alpha=0.7):
        """Create a colormap that combines base image with tumor overlay"""
        from matplotlib.colors import to_rgba
        
        # Get the tumor color as RGBA
        tumor_rgba = to_rgba(tumor_color, alpha=tumor_alpha)
        
        return {
            'base_cmap': base_colormap,
            'tumor_color': tumor_rgba,
            'tumor_alpha': tumor_alpha
        }
    
    def create_mosaic_view(self, nifti_file, rows=3, cols=5, slice_gap=2, 
                          contrast='T1', window_level=None, window_width=None,
                          colormap='gray', timepoint_info="", segmentation_file=None,
                          tumor_color='red', tumor_alpha=0.5, show_contour=True,
                          contour_width=2, contour_style='solid', tumor_volume_info=None):
        """
        Create a mosaic view of brain slices with optional tumor segmentation overlay
        
        Parameters:
        - nifti_file: Path to NIFTI file
        - rows, cols: Mosaic grid dimensions
        - slice_gap: Gap between slices
        - contrast: Image contrast type (for labeling)
        - window_level, window_width: Windowing parameters
        - colormap: Matplotlib colormap for brain image
        - timepoint_info: Additional info for title
        - segmentation_file: Path to tumor segmentation file
        - tumor_color: Color for tumor overlay ('red', 'yellow', 'cyan', etc.)
        - tumor_alpha: Transparency of tumor overlay (0-1)
        - show_contour: Whether to show tumor contour instead of filled overlay
        - contour_width: Width of contour lines
        - contour_style: Style of contour lines ('solid', 'dashed', 'dotted')
        - tumor_volume_info: Dictionary with tumor volume information
        """
        if not nifti_file or not os.path.exists(nifti_file):
            print(f"Error: File not found: {nifti_file}")
            return None
            
        # Load NIFTI file
        img = nib.load(nifti_file)
        data = img.get_fdata()
        
        # Load segmentation if provided
        seg_data = None
        if segmentation_file and os.path.exists(segmentation_file):
            seg_img = nib.load(segmentation_file)
            seg_data = seg_img.get_fdata()
            print(f"Loaded segmentation: {segmentation_file}")
        
        # Handle different orientations - use axial slices (last dimension)
        if len(data.shape) == 4:
            data = data[:, :, :, 0]  # Take first volume if 4D
        
        if seg_data is not None and len(seg_data.shape) == 4:
            seg_data = seg_data[:, :, :, 0]  # Take first volume if 4D
        
        # Get slice indices - prioritize tumor-containing slices if segmentation is available
        total_slices = data.shape[2]
        num_slices = rows * cols
        
        if seg_data is not None:
            # Find slices with tumor
            tumor_slices = []
            for z in range(total_slices):
                slice_volume = np.sum(seg_data[:, :, z] > 0.5)
                if slice_volume > 25:  # Minimum 25 voxels to be considered a tumor slice
                    tumor_slices.append(z)
            
            if len(tumor_slices) >= num_slices:
                # If we have enough tumor slices, select evenly from them
                slice_indices = np.linspace(0, len(tumor_slices)-1, num_slices, dtype=int)
                slice_indices = [tumor_slices[i] for i in slice_indices]
                print(f"Selected {num_slices} tumor-containing slices from {len(tumor_slices)} available")
            elif len(tumor_slices) > 0:
                # If we have some tumor slices but not enough, use all tumor slices and fill with adjacent slices
                tumor_range_start = max(0, min(tumor_slices) - 5)
                tumor_range_end = min(total_slices, max(tumor_slices) + 6)
                extended_range = list(range(tumor_range_start, tumor_range_end))
                slice_indices = np.linspace(0, len(extended_range)-1, num_slices, dtype=int)
                slice_indices = [extended_range[i] for i in slice_indices]
                print(f"Using {len(tumor_slices)} tumor slices, extended range to get {num_slices} total slices")
            else:
                # Fallback to central slices if no tumor found
                start_idx = int(total_slices * 0.2)
                end_idx = int(total_slices * 0.8)
                slice_indices = np.linspace(start_idx, end_idx, num_slices, dtype=int)
                print("No tumor slices found, using central brain slices")
        else:
            # No segmentation available - use central slices
            start_idx = int(total_slices * 0.2)
            end_idx = int(total_slices * 0.8)
            slice_indices = np.linspace(start_idx, end_idx, num_slices, dtype=int)
        
        # Create figure
        fig = plt.figure(figsize=(cols * 3, rows * 3))
        fig.patch.set_facecolor('black')  # Make figure background black

        gs = GridSpec(rows, cols, figure=fig, hspace=0.1, wspace=0.1)
        
        # Set windowing if provided
        if window_level is not None and window_width is not None:
            vmin = window_level - window_width / 2
            vmax = window_level + window_width / 2
        else:
            # Auto-windowing
            vmin, vmax = np.percentile(data[data > 0], [2, 98])
        
        for i, slice_idx in enumerate(slice_indices):
            row = i // cols
            col = i % cols
            
            ax = fig.add_subplot(gs[row, col])
            ax.set_facecolor('black')  # Make axes (subplot) background black
            
            # Get slice in proper neurological orientation (LPS)
            # For axial slices: transpose and flip for radiological convention
            slice_data = data[:, :, slice_idx].T
            # Remove the flipud to fix upside-down issue
            # slice_data = np.flipud(slice_data)  # Commented out to fix orientation
            
            # Display brain slice
            im = ax.imshow(slice_data, cmap=colormap, vmin=vmin, vmax=vmax, 
                          aspect='equal', interpolation='nearest')
            
            # Add tumor overlay if segmentation is available
            if seg_data is not None:
                seg_slice = seg_data[:, :, slice_idx].T
                # seg_slice = np.flipud(seg_slice)  # Commented out to match brain orientation
                
                # Much more robust tumor detection
                # First, check the actual range of values in this slice
                slice_max = np.max(seg_slice)
                slice_min = np.min(seg_slice[seg_slice > 0]) if np.any(seg_slice > 0) else 0
                
                # Adaptive threshold based on actual data values
                if slice_max > 0.8:
                    tumor_threshold = 0.5  # Standard threshold for binary masks
                elif slice_max > 0.1:
                    tumor_threshold = slice_max * 0.5  # Use 50% of max value
                else:
                    tumor_threshold = 0.05  # Very low threshold for normalized data
                
                tumor_mask = seg_slice > tumor_threshold
                
                # Require substantial tumor area (at least 25 pixels) and check that
                # the tumor region actually corresponds to visible brain tissue
                if np.sum(tumor_mask) > 25:
                    # Additional check: make sure tumor is in a region with brain tissue
                    # (avoid artifacts in background regions)
                    brain_slice = slice_data
                    brain_threshold = np.percentile(brain_slice[brain_slice > 0], 10) if np.any(brain_slice > 0) else 0
                    
                    # Only show tumor if it overlaps with brain tissue
                    brain_mask = brain_slice > brain_threshold
                    tumor_in_brain = tumor_mask & brain_mask
                    
                    if np.sum(tumor_in_brain) > 15:  # At least 15 pixels of tumor in brain tissue
                        if show_contour:
                            # Show tumor contour with adaptive threshold and line style
                            linestyles = contour_style if contour_style in ['solid', 'dashed', 'dotted', 'dashdot'] else 'solid'
                            contours = ax.contour(seg_slice, levels=[tumor_threshold], colors=[tumor_color], 
                                                linewidths=contour_width, linestyles=linestyles, alpha=tumor_alpha)
                        else:
                            # Show filled tumor overlay only where there's brain tissue
                            tumor_overlay = np.zeros((*tumor_in_brain.shape, 4))  # RGBA
                            from matplotlib.colors import to_rgba
                            tumor_rgba = to_rgba(tumor_color, alpha=tumor_alpha)
                            tumor_overlay[tumor_in_brain] = tumor_rgba
                            ax.imshow(tumor_overlay, aspect='equal', interpolation='nearest')
            
            ax.set_xticks([])
            ax.set_yticks([])
            ax.set_title(f'Slice {slice_idx}', fontsize=8, color='white')
            
            # Add LPS orientation annotations (only on first slice for clarity)
            if i == 0:
                # Left-Right annotation (horizontal)
                ax.text(0.02, 0.5, 'L', transform=ax.transAxes, 
                       fontsize=12, color='yellow', ha='left', va='center', fontweight='bold')
                ax.text(0.98, 0.5, 'R', transform=ax.transAxes, 
                       fontsize=12, color='yellow', ha='right', va='center', fontweight='bold')
                
                # Posterior-Anterior annotation (vertical) 
                ax.text(0.5, 0.02, 'P', transform=ax.transAxes, 
                       fontsize=12, color='yellow', ha='center', va='bottom', fontweight='bold')
                ax.text(0.5, 0.98, 'A', transform=ax.transAxes, 
                       fontsize=12, color='yellow', ha='center', va='top', fontweight='bold')
        
        # Add title with timepoint info, segmentation status, and tumor volume
        seg_status = " (with tumor overlay)" if seg_data is not None else ""
        volume_info = ""
        if tumor_volume_info:
            volume_info = f" | Volume: {tumor_volume_info['volume_ml']:.2f} mL"
        
        title = f'{contrast} - {timepoint_info}{seg_status}{volume_info}'
        fig.suptitle(title, fontsize=16, fontweight='bold', color='white')
        
        plt.tight_layout()
        return fig
    
    def create_flipbook_slides(self, contrast_type, rows=3, cols=5, 
                              window_level=None, window_width=None,
                              colormap='gray', tumor_color='red', 
                              tumor_alpha=0.5, show_contour=True,
                              contour_width=2, contour_style='solid'):
        """Create individual slides for the flipbook with tumor overlay"""
        if not self.timepoint_folders:
            self.find_timepoint_folders()
            
        if not self.contrast_types:
            self.find_contrast_types()
            
        if contrast_type not in self.contrast_types:
            print(f"Error: Contrast type '{contrast_type}' not found.")
            print(f"Available types: {self.contrast_types}")
            return []
        
        slide_files = []
        volume_data = []  # Store volume information for all timepoints
        
        for i, tp_folder in enumerate(self.timepoint_folders):
            timepoint_name = tp_folder['folder_name']
            print(f"Creating slide {i+1}/{len(self.timepoint_folders)}: {timepoint_name}")
            
            # Get the NIFTI file for this timepoint and contrast
            nifti_file = self.get_file_for_timepoint_and_contrast(tp_folder, contrast_type)
            
            if not nifti_file:
                print(f"Skipping timepoint {timepoint_name} - no {contrast_type} file found")
                continue
            
            # Get segmentation file for this timepoint
            seg_file = self.get_segmentation_file(tp_folder)
            
            # Calculate tumor volume
            tumor_volume_info = None
            if seg_file:
                tumor_volume_info = self.calculate_tumor_volume(seg_file)
                if tumor_volume_info:
                    volume_data.append({
                        'timepoint': timepoint_name,
                        'volume_ml': tumor_volume_info['volume_ml'],
                        'volume_mm3': tumor_volume_info['volume_mm3'],
                        'voxel_count': tumor_volume_info['voxel_count']
                    })
                    print(f"  Tumor volume: {tumor_volume_info['volume_ml']:.2f} mL")
                
            # Create mosaic view with tumor overlay and volume info
            fig = self.create_mosaic_view(
                nifti_file, rows=rows, cols=cols, contrast=contrast_type,
                window_level=window_level, window_width=window_width,
                colormap=colormap, timepoint_info=f"Timepoint {timepoint_name}",
                segmentation_file=seg_file, tumor_color=tumor_color,
                tumor_alpha=tumor_alpha, show_contour=show_contour,
                contour_width=contour_width, contour_style=contour_style, tumor_volume_info=tumor_volume_info
            )
            
            if fig is None:
                continue
                
            # Save slide
            slide_file = os.path.join(self.output_folder, f'slide_{i+1:03d}_{contrast_type}_{timepoint_name}.png')
            
            # Force overwrite by removing existing file
            if os.path.exists(slide_file):
                os.remove(slide_file)
                print(f"Removed existing file: {slide_file}")
            
            plt.savefig(slide_file, dpi=150, bbox_inches='tight', facecolor='black', edgecolor='none')
            plt.close(fig)
            
            print(f"Saved new slide: {slide_file}")
            slide_files.append(slide_file)
        
        # Print volume progression summary and create volume plot
        if volume_data:
            print(f"\n=== TUMOR VOLUME PROGRESSION ({contrast_type}) ===")
            for vol_info in volume_data:
                print(f"Timepoint {vol_info['timepoint']}: {vol_info['volume_ml']:.2f} mL")
            
            # Calculate volume change
            if len(volume_data) > 1:
                initial_vol = volume_data[0]['volume_ml']
                final_vol = volume_data[-1]['volume_ml']
                vol_change = final_vol - initial_vol
                vol_percent_change = (vol_change / initial_vol) * 100 if initial_vol > 0 else 0
                print(f"Volume change: {vol_change:+.2f} mL ({vol_percent_change:+.1f}%)")
                
                # Create volume progression plot
                volume_plot_file = self.create_volume_progression_plot(volume_data, contrast_type)
        
        return slide_files
    
    def create_volume_progression_plot(self, volume_data, contrast_type='T1CE', output_folder=None):
        """
        Create a tumor volume progression plot across timepoints
        
        Parameters:
        - volume_data: List of dictionaries with 'timepoint' and 'volume_ml' keys
        - contrast_type: Type of contrast for labeling
        - output_folder: Folder to save the plot (optional)
        
        Returns:
        - Path to saved plot file
        """
        if not volume_data or len(volume_data) < 2:
            print("Warning: Need at least 2 timepoints with volume data to create progression plot")
            return None
        
        # Extract timepoints and volumes
        timepoints = [vol_info['timepoint'] for vol_info in volume_data]
        volumes = [vol_info['volume_ml'] for vol_info in volume_data]
        
        # Create figure with dark theme
        plt.style.use('dark_background')
        fig, ax = plt.subplots(figsize=(10, 6))
        fig.patch.set_facecolor('black')
        ax.set_facecolor('black')
        
        # Plot volume progression
        ax.plot(timepoints, volumes, 'o-', linewidth=3, markersize=8, 
                color='#FF6B6B', markerfacecolor='#FF4444', markeredgecolor='white', markeredgewidth=2)
        
        # Fill area under curve for visual appeal
        ax.fill_between(timepoints, volumes, alpha=0.3, color='#FF6B6B')
        
        # Customize plot
        ax.set_xlabel('Timepoint', fontsize=14, fontweight='bold', color='white')
        ax.set_ylabel('Tumor Volume (mL)', fontsize=14, fontweight='bold', color='white')
        ax.set_title(f'Tumor Volume Progression - {contrast_type}', 
                    fontsize=16, fontweight='bold', color='white', pad=20)
        
        # Add volume values as labels on points
        for i, (tp, vol) in enumerate(zip(timepoints, volumes)):
            ax.annotate(f'{vol:.1f} mL', 
                       (tp, vol), 
                       textcoords="offset points", 
                       xytext=(0,15), 
                       ha='center', 
                       fontsize=10, 
                       color='white',
                       fontweight='bold',
                       bbox=dict(boxstyle="round,pad=0.3", facecolor='black', alpha=0.7))
        
        # Calculate and display volume change statistics
        initial_vol = volumes[0]
        final_vol = volumes[-1]
        vol_change = final_vol - initial_vol
        vol_percent_change = (vol_change / initial_vol) * 100 if initial_vol > 0 else 0
        
        # Add statistics text box
        stats_text = f'Initial: {initial_vol:.1f} mL\nFinal: {final_vol:.1f} mL\nChange: {vol_change:+.1f} mL ({vol_percent_change:+.1f}%)'
        ax.text(0.02, 0.98, stats_text, transform=ax.transAxes, 
                fontsize=11, verticalalignment='top',
                bbox=dict(boxstyle="round,pad=0.5", facecolor='#333333', alpha=0.8),
                color='white', fontweight='bold')
        
        # Style the grid and spines
        ax.grid(True, alpha=0.3, color='gray', linestyle='--')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('white')
        ax.spines['bottom'].set_color('white')
        ax.tick_params(colors='white')
        
        # Set reasonable y-axis limits with some padding
        y_min, y_max = min(volumes), max(volumes)
        y_range = y_max - y_min
        if y_range > 0:
            ax.set_ylim(max(0, y_min - y_range * 0.1), y_max + y_range * 0.1)
        
        # Ensure x-axis shows all timepoints
        ax.set_xticks(timepoints)
        
        plt.tight_layout()
        
        # Save plot
        if output_folder is None:
            output_folder = self.output_folder
        
        plot_file = os.path.join(output_folder, f"tumor_volume_progression_{contrast_type}.png")
        plt.savefig(plot_file, dpi=300, bbox_inches='tight', facecolor='black', edgecolor='none')
        plt.close(fig)
        
        # Reset matplotlib style to default
        plt.style.use('default')
        
        print(f"✅ Saved volume progression plot: {plot_file}")
        return plot_file
    
    def create_tumor_focused_summary_slide(self, contrast_type, num_slices=4,
                                         window_level=None, window_width=None,
                                         colormap='gray', tumor_color='red',
                                         tumor_alpha=0.5, show_contour=True,
                                         contour_width=2, contour_style='solid'):
        """
        Create a summary slide showing tumor progression across timepoints
        
        Shows 3-4 central tumor slices as rows, with each timepoint as a column.
        This provides a focused view of tumor evolution over time.
        
        Parameters:
        - contrast_type: Type of contrast (T1, T2, T1CE, etc.)
        - num_slices: Number of central tumor slices to show (3-4 recommended)
        - Other parameters: Same as create_mosaic_view
        
        Returns:
        - Path to saved summary slide image
        """
        if not self.timepoint_folders:
            self.find_timepoint_folders()
            
        if not self.contrast_types:
            self.find_contrast_types()
            
        if contrast_type not in self.contrast_types:
            print(f"Error: Contrast type '{contrast_type}' not found.")
            return None
            
        print(f"Creating tumor-focused summary slide for {contrast_type}...")
        
        # Collect all tumor data and find central slices
        timepoint_data = []
        all_tumor_slices = set()
        
        for tp_folder in self.timepoint_folders:
            timepoint_name = tp_folder['folder_name']
            
            # Get the NIFTI file for this timepoint and contrast
            nifti_file = self.get_file_for_timepoint_and_contrast(tp_folder, contrast_type)
            if not nifti_file:
                continue
                
            # Get segmentation file
            seg_file = self.get_segmentation_file(tp_folder)
            if not seg_file:
                continue
                
            # Load segmentation to find tumor slices
            try:
                seg_img = nib.load(seg_file)
                seg_data = seg_img.get_fdata()
                
                # Find slices with tumor
                tumor_slices = []
                for z in range(seg_data.shape[2]):
                    if np.any(seg_data[:, :, z] > 0):
                        tumor_slices.append(z)
                        all_tumor_slices.add(z)
                
                if tumor_slices:
                    timepoint_data.append({
                        'name': timepoint_name,
                        'nifti_file': nifti_file,
                        'seg_file': seg_file,
                        'tumor_slices': tumor_slices,
                        'folder': tp_folder
                    })
                    
            except Exception as e:
                print(f"Warning: Could not process segmentation for {timepoint_name}: {e}")
                continue
        
        if not timepoint_data or not all_tumor_slices:
            print("No tumor data found for summary slide")
            return None
            
        # Select central tumor slices
        sorted_tumor_slices = sorted(list(all_tumor_slices))
        if len(sorted_tumor_slices) >= num_slices:
            # Take evenly spaced slices from the middle portion
            start_idx = len(sorted_tumor_slices) // 4
            end_idx = 3 * len(sorted_tumor_slices) // 4
            middle_slices = sorted_tumor_slices[start_idx:end_idx]
            
            if len(middle_slices) >= num_slices:
                step = len(middle_slices) // num_slices
                selected_slices = [middle_slices[i * step] for i in range(num_slices)]
            else:
                selected_slices = middle_slices
        else:
            selected_slices = sorted_tumor_slices
            
        selected_slices = selected_slices[:num_slices]  # Ensure we don't exceed num_slices
        
        print(f"Selected tumor slices: {selected_slices}")
        
        # Create the summary figure
        # Rows = slices, Columns = timepoints (as requested)
        num_timepoints = len(timepoint_data)
        fig_width = num_timepoints * 4  # 4 inches per timepoint column
        fig_height = num_slices * 3     # 3 inches per slice row
        
        fig, axes = plt.subplots(num_slices, num_timepoints, 
                                figsize=(fig_width, fig_height), 
                                facecolor='black')
        
        # Handle single row or column cases
        if num_slices == 1:
            axes = axes.reshape(1, -1)
        elif num_timepoints == 1:
            axes = axes.reshape(-1, 1)
            
        fig.suptitle(f'Tumor Evolution Summary - {contrast_type}', 
                     fontsize=20, color='white', y=0.92)
        
        # Process each slice (row) and timepoint (column)
        for row, slice_idx in enumerate(selected_slices):
            for col, tp_data in enumerate(timepoint_data):
                # Load brain image for this timepoint
                img = nib.load(tp_data['nifti_file'])
                data = img.get_fdata()
                
                # Load segmentation for this timepoint
                seg_img = nib.load(tp_data['seg_file'])
                seg_data = seg_img.get_fdata()
                
                # Calculate tumor volume for this timepoint (for labeling)
                tumor_volume_info = self.calculate_tumor_volume(tp_data['seg_file'])
                volume_text = ""
                if tumor_volume_info:
                    volume_text = f" ({tumor_volume_info['volume_ml']:.1f} mL)"
                ax = axes[row, col]
                ax.set_facecolor('black')
                
                # Get brain slice
                slice_data = data[:, :, slice_idx].T
                
                # Set windowing - use the same approach as create_mosaic_view
                if window_level is not None and window_width is not None:
                    vmin = window_level - window_width / 2
                    vmax = window_level + window_width / 2
                else:
                    # Auto-windowing using percentile approach
                    vmin, vmax = np.percentile(data[data > 0], [2, 98])
                
                # Display brain slice
                im = ax.imshow(slice_data, cmap=colormap, vmin=vmin, vmax=vmax, 
                              aspect='equal', origin='lower')
                
                # Add tumor overlay
                seg_slice = seg_data[:, :, slice_idx].T
                if np.any(seg_slice > 0):
                    if show_contour:
                        # Create contour overlay - use same approach as create_mosaic_view
                        tumor_threshold = 0.5
                        linestyles = contour_style if contour_style in ['solid', 'dashed', 'dotted', 'dashdot'] else 'solid'
                        contours = ax.contour(seg_slice, levels=[tumor_threshold], colors=[tumor_color], 
                                            linewidths=contour_width, linestyles=linestyles, alpha=tumor_alpha)
                    else:
                        # Create filled overlay - use same approach as create_mosaic_view  
                        from matplotlib.colors import to_rgba
                        tumor_mask = seg_slice > 0
                        tumor_overlay = np.zeros((*tumor_mask.shape, 4))  # RGBA
                        tumor_rgba = to_rgba(tumor_color, alpha=tumor_alpha)
                        tumor_overlay[tumor_mask] = tumor_rgba
                        ax.imshow(tumor_overlay, aspect='equal', origin='lower', interpolation='nearest')
                
                # Remove ticks and add labels
                ax.set_xticks([])
                ax.set_yticks([])
                
                # Add timepoint label on top row (columns are timepoints)
                if row == 0:
                    ax.text(0.5, 1.02, f'{tp_data["name"]}{volume_text}', 
                           transform=ax.transAxes, fontsize=10, color='white', 
                           ha='center', va='bottom', fontweight='bold')
                
                # Add slice number on first column (rows are slices)
                if col == 0:
                    ax.text(-0.1, 0.5, f'Slice {slice_idx}', 
                           transform=ax.transAxes, rotation=90,
                           fontsize=12, color='white', ha='right', va='center',
                           fontweight='bold')
        
        # Adjust layout with more space for title and labels
        plt.tight_layout()
        plt.subplots_adjust(top=0.88, hspace=0.12, wspace=0.05)
        
        # Save the summary slide
        summary_file = os.path.join(self.output_folder, f"tumor_summary_{contrast_type}.png")
        plt.savefig(summary_file, dpi=300, bbox_inches='tight', 
                   facecolor='black', edgecolor='none')
        plt.close(fig)
        
        print(f"✅ Created tumor-focused summary slide: {summary_file}")
        return summary_file
    
    def create_animated_gif(self, slide_files, output_path, duration=800, loop=0):
        """
        Create an animated GIF from flipbook slides
        
        Parameters:
        - slide_files: List of paths to slide image files
        - output_path: Path for output GIF file
        - duration: Duration of each frame in milliseconds (default 800ms = 0.8s)
        - loop: Number of loops (0 = infinite loop)
        
        Returns:
        - Path to created GIF file or None if failed
        """
        if not slide_files:
            print("No slide files provided for GIF creation")
            return None
            
        try:
            from PIL import Image
            
            # Load all images
            images = []
            for slide_file in slide_files:
                if os.path.exists(slide_file):
                    img = Image.open(slide_file)
                    # Convert to RGB if necessary (for GIF compatibility)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    images.append(img)
                else:
                    print(f"Warning: Slide file not found: {slide_file}")
            
            if not images:
                print("No valid images found for GIF creation")
                return None
            
            # Create animated GIF
            print(f"Creating animated GIF with {len(images)} frames...")
            images[0].save(
                output_path,
                save_all=True,
                append_images=images[1:],
                duration=duration,
                loop=loop,
                optimize=True
            )
            
            print(f"✅ Created animated GIF: {output_path}")
            return output_path
            
        except ImportError:
            print("Error: PIL/Pillow not available for GIF creation")
            return None
        except Exception as e:
            print(f"Error creating GIF: {e}")
            return None
    
    def create_enhanced_animated_gif(self, slide_files, output_path, duration=1000, 
                                   loop=0, add_progress_bar=True, add_frame_numbers=True):
        """
        Create an enhanced animated GIF with optional progress indicators
        
        Parameters:
        - slide_files: List of paths to slide image files
        - output_path: Path for output GIF file
        - duration: Duration of each frame in milliseconds
        - loop: Number of loops (0 = infinite)
        - add_progress_bar: Add a progress bar at the bottom
        - add_frame_numbers: Add frame numbers to each slide
        
        Returns:
        - Path to created GIF file or None if failed
        """
        if not slide_files:
            print("No slide files provided for GIF creation")
            return None
            
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            # Load and process all images
            processed_images = []
            total_frames = len(slide_files)
            
            for i, slide_file in enumerate(slide_files):
                if not os.path.exists(slide_file):
                    print(f"Warning: Slide file not found: {slide_file}")
                    continue
                    
                img = Image.open(slide_file)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Create a copy to modify
                enhanced_img = img.copy()
                draw = ImageDraw.Draw(enhanced_img)
                
                # Add frame number if requested
                if add_frame_numbers:
                    frame_text = f"{i+1}/{total_frames}"
                    try:
                        # Try to use a larger font
                        font = ImageFont.truetype("arial.ttf", 24)
                    except:
                        # Fall back to default font
                        font = ImageFont.load_default()
                    
                    # Add frame number in top-right corner
                    text_bbox = draw.textbbox((0, 0), frame_text, font=font)
                    text_width = text_bbox[2] - text_bbox[0]
                    text_height = text_bbox[3] - text_bbox[1]
                    
                    x = enhanced_img.width - text_width - 20
                    y = 20
                    
                    # Add background rectangle for better visibility
                    draw.rectangle([x-5, y-5, x+text_width+5, y+text_height+5], 
                                 fill=(0, 0, 0, 128))
                    draw.text((x, y), frame_text, fill=(255, 255, 255), font=font)
                
                # Add progress bar if requested
                if add_progress_bar:
                    bar_height = 8
                    bar_y = enhanced_img.height - bar_height - 10
                    bar_width = enhanced_img.width - 40
                    bar_x = 20
                    
                    # Background bar
                    draw.rectangle([bar_x, bar_y, bar_x + bar_width, bar_y + bar_height], 
                                 fill=(50, 50, 50))
                    
                    # Progress fill
                    progress = (i + 1) / total_frames
                    fill_width = int(bar_width * progress)
                    draw.rectangle([bar_x, bar_y, bar_x + fill_width, bar_y + bar_height], 
                                 fill=(0, 150, 255))
                
                processed_images.append(enhanced_img)
            
            if not processed_images:
                print("No valid images found for GIF creation")
                return None
            
            # Create animated GIF
            print(f"Creating enhanced animated GIF with {len(processed_images)} frames...")
            processed_images[0].save(
                output_path,
                save_all=True,
                append_images=processed_images[1:],
                duration=duration,
                loop=loop,
                optimize=True
            )
            
            print(f"✅ Created enhanced animated GIF: {output_path}")
            return output_path
            
        except ImportError:
            print("Error: PIL/Pillow not available for GIF creation")
            return None
        except Exception as e:
            print(f"Error creating enhanced GIF: {e}")
            return None
    
    
    def generate_flipbook(self, contrast_type, rows=3, cols=5, 
                         window_level=None, window_width=None,
                         colormap='gray', title=None, tumor_color='red',
                         tumor_alpha=0.5, show_contour=True, contour_width=2, contour_style='solid',
                         create_summary_slide=True, create_animated_gif=True, 
                         gif_duration=1000, summary_slices=4):
        """
        Generate complete flipbook for a specific contrast type with tumor overlay
        
        Parameters:
        - contrast_type: Type of contrast (T1, T2, T1CE, FL, etc.)
        - rows, cols: Mosaic grid dimensions
        - window_level, window_width: Windowing parameters
        - colormap: Matplotlib colormap for brain image
        - title: Flipbook title
        - tumor_color: Color for tumor overlay ('red', 'yellow', 'cyan', 'green', 'magenta')
        - tumor_alpha: Transparency of tumor overlay (0-1)
        - show_contour: If True, show tumor contour; if False, show filled overlay
        - contour_width: Width of contour lines
        - contour_style: Style of contour lines ('solid', 'dashed', 'dotted')
        - create_summary_slide: If True, create tumor-focused summary slide (default: True)
        - create_animated_gif: If True, create animated GIF (default: True) 
        - gif_duration: Duration per frame in animated GIF in milliseconds (default: 1000)
        - summary_slices: Number of central tumor slices in summary slide (default: 4)
        """
        if title is None:
            seg_suffix = " with Tumor Overlay" if self.segmentation_folder else ""
            title = f"Brain Tumor Flipbook - {contrast_type}{seg_suffix}"
            
        print("=== Enhanced Brain Tumor Flipbook Generator ===")
        print(f"Input folder: {self.registered_folder}")
        print(f"Segmentation folder: {self.segmentation_folder}")
        print(f"Output folder: {self.output_folder}")
        print(f"Contrast type: {contrast_type}")
        print(f"Tumor visualization: {'Contour' if show_contour else 'Filled overlay'} in {tumor_color}")
        
        # Find timepoint folders and contrast types
        self.find_timepoint_folders()
        self.find_contrast_types()
        
        # Create slides with tumor overlay
        slide_files = self.create_flipbook_slides(
            contrast_type=contrast_type, rows=rows, cols=cols,
            window_level=window_level, window_width=window_width,
            colormap=colormap, tumor_color=tumor_color,
            tumor_alpha=tumor_alpha, show_contour=show_contour,
            contour_width=contour_width, contour_style=contour_style
        )
        
        if not slide_files:
            print("Error: No slides were created. Check your contrast type and files.")
            return []
        
        print(f"\nFlipbook generated successfully!")
        print(f"- {len(slide_files)} slides created")
        print(f"- Individual slides saved in: {self.output_folder}")
        
        # Create tumor-focused summary slide (the radiologist-recommended feature)
        summary_slide = None
        if create_summary_slide and self.segmentation_folder:
            print(f"\n=== Creating Tumor-Focused Summary Slide ===")
            summary_slide = self.create_tumor_focused_summary_slide(
                contrast_type=contrast_type,
                num_slices=summary_slices,
                window_level=window_level,
                window_width=window_width,
                colormap=colormap,
                tumor_color=tumor_color,
                tumor_alpha=tumor_alpha,
                show_contour=show_contour,
                contour_width=contour_width,
                contour_style=contour_style
            )
        elif create_summary_slide and not self.segmentation_folder:
            print("Note: Skipping summary slide creation (no segmentation data)")
        
        # Create animated GIF from the original flipbook slides (excluding summary)
        animated_gif = None
        if create_animated_gif:
            print(f"\n=== Creating Animated GIF ===")
            gif_path = os.path.join(self.output_folder, f"flipbook_animation_{contrast_type}.gif")
            
            # Use enhanced GIF creation with progress indicators
            animated_gif = self.create_enhanced_animated_gif(
                slide_files=slide_files,  # Only use original flipbook slides, not summary
                output_path=gif_path,
                duration=gif_duration,
                loop=0,  # Infinite loop
                add_progress_bar=True,
                add_frame_numbers=True
            )
        
        # Prepare return results
        results = {
            'slide_files': slide_files,
            'slide_count': len(slide_files),
            'summary_slide': summary_slide,
            'animated_gif': animated_gif,
            'output_folder': self.output_folder
        }
        
        # Print summary
        print(f"\n✅ Complete flipbook package created:")
        print(f"   - {len(slide_files)} individual slides")
        if summary_slide:
            print(f"   - Tumor evolution summary slide")
        if animated_gif:
            print(f"   - Animated GIF visualization")
        print(f"   - All files saved in: {self.output_folder}")
        
        return results
    
    def generate_all_contrast_flipbooks(self, rows=3, cols=5, 
                                       window_level=None, window_width=None,
                                       colormap='gray', tumor_color='red',
                                       tumor_alpha=0.5, show_contour=True, 
                                       contour_width=2, contour_style='solid',
                                       create_summary_slide=True, create_animated_gif=True,
                                       gif_duration=1000, summary_slices=4):
        """Generate flipbooks for all available contrast types with tumor overlay"""
        self.find_timepoint_folders()
        self.find_contrast_types()
        
        print("=== Generating Flipbooks for All Contrast Types ===")
        print(f"Found contrast types: {self.contrast_types}")
        print(f"Tumor visualization: {'Contour' if show_contour else 'Filled overlay'} in {tumor_color}")
        
        results = {}
        
        for i, contrast in enumerate(self.contrast_types):
            print(f"\n=== Generating flipbook {i+1}/{len(self.contrast_types)} for {contrast} ===")
            
            # Create contrast-specific output folder
            contrast_output_folder = os.path.join(self.output_folder, f"{contrast}_flipbook")
            os.makedirs(contrast_output_folder, exist_ok=True)
            
            # Temporarily change the output folder for this contrast
            original_output = self.output_folder
            self.output_folder = contrast_output_folder
            
            flipbook_results = self.generate_flipbook(
                contrast_type=contrast, rows=rows, cols=cols,
                window_level=window_level, window_width=window_width,
                colormap=colormap, tumor_color=tumor_color,
                tumor_alpha=tumor_alpha, show_contour=show_contour,
                contour_width=contour_width, contour_style=contour_style,
                create_summary_slide=create_summary_slide, 
                create_animated_gif=create_animated_gif,
                gif_duration=gif_duration, summary_slices=summary_slices
            )
            
            # Restore original output folder
            self.output_folder = original_output
            
            # Handle both old (list) and new (dict) return formats
            if isinstance(flipbook_results, dict):
                results[contrast] = flipbook_results
                results[contrast]['output_folder'] = contrast_output_folder
            else:
                # Legacy format (list of slide files)
                results[contrast] = {
                    'slide_files': flipbook_results,
                    'slide_count': len(flipbook_results) if flipbook_results else 0,
                    'output_folder': contrast_output_folder
                }
        
        print(f"\n=== ALL FLIPBOOKS GENERATED SUCCESSFULLY ===")
        print(f"Generated {len(results)} flipbooks for contrast types: {list(results.keys())}")
        print(f"\n*** FILES SAVED TO: {os.path.abspath(self.output_folder)} ***")
        
        # Print absolute paths for verification
        for contrast, result in results.items():
            slide_count = result.get('slide_count', len(result.get('slide_files', [])))
            if slide_count > 0:
                extras = []
                if result.get('summary_slide'):
                    extras.append("summary slide")
                if result.get('animated_gif'):
                    extras.append("animated GIF")
                
                extras_text = f" + {', '.join(extras)}" if extras else ""
                print(f"  {contrast}: {slide_count} slides{extras_text} in {os.path.abspath(result['output_folder'])}")
        
        return results
    
    def generate_all_contrast_flipbooks_slides_only(self, rows=3, cols=5, 
                                                   window_level=None, window_width=None,
                                                   colormap='gray', tumor_color='red',
                                                   tumor_alpha=0.5, show_contour=True, 
                                                   contour_width=2, contour_style='solid', create_pptx=True,
                                                   create_summary_slide=True, create_animated_gif=True,
                                                   gif_duration=1000, summary_slices=4):
        """Generate flipbook slides for all available contrast types with optional PPTX output"""
        self.find_timepoint_folders()
        self.find_contrast_types()
        
        print("=== Generating Flipbook Slides for All Contrast Types ===")
        print(f"Found contrast types: {self.contrast_types}")
        print(f"Tumor visualization: {'Contour' if show_contour else 'Filled overlay'} in {tumor_color}")
        print(f"PowerPoint output: {'Enabled' if create_pptx and PPTX_AVAILABLE else 'Disabled'}")
        
        results = {}
        
        for i, contrast in enumerate(self.contrast_types):
            print(f"\n=== Generating slides {i+1}/{len(self.contrast_types)} for {contrast} ===")
            
            # Create contrast-specific output folder
            contrast_output_folder = os.path.join(self.output_folder, f"{contrast}_flipbook")
            os.makedirs(contrast_output_folder, exist_ok=True)
            
            # Temporarily change the output folder for this contrast
            original_output = self.output_folder
            self.output_folder = contrast_output_folder
            
            # Generate slides
            slide_files = self.create_flipbook_slides(
                contrast_type=contrast, rows=rows, cols=cols,
                window_level=window_level, window_width=window_width,
                colormap=colormap, tumor_color=tumor_color,
                tumor_alpha=tumor_alpha, show_contour=show_contour,
                contour_width=contour_width, contour_style=contour_style
            )
            
            # Create tumor-focused summary slide if requested
            summary_slide = None
            if create_summary_slide and self.segmentation_folder and slide_files:
                print(f"  Creating tumor-focused summary slide for {contrast}...")
                summary_slide = self.create_tumor_focused_summary_slide(
                    contrast_type=contrast,
                    num_slices=summary_slices,
                    window_level=window_level,
                    window_width=window_width,
                    colormap=colormap,
                    tumor_color=tumor_color,
                    tumor_alpha=tumor_alpha,
                    show_contour=show_contour,
                    contour_width=contour_width,
                    contour_style=contour_style
                )
            elif create_summary_slide and not self.segmentation_folder:
                print(f"  Note: Skipping summary slide for {contrast} (no segmentation data)")
            
            # Create animated GIF if requested
            animated_gif = None
            if create_animated_gif and slide_files:
                print(f"  Creating animated GIF for {contrast}...")
                gif_path = os.path.join(contrast_output_folder, f"flipbook_animation_{contrast}.gif")
                animated_gif = self.create_enhanced_animated_gif(
                    slide_files=slide_files,
                    output_path=gif_path,
                    duration=gif_duration,
                    loop=0,
                    add_progress_bar=True,
                    add_frame_numbers=True
                )
            
            # Generate PowerPoint if requested and available
            pptx_file = None
            if create_pptx and PPTX_AVAILABLE and slide_files:
                seg_suffix = " with Tumor Overlay" if self.segmentation_folder else ""
                title = f"Brain Tumor Flipbook - {contrast}{seg_suffix}"
                pptx_file = self.create_powerpoint_flipbook(slide_files, contrast, title=title)
            
            # Restore original output folder
            self.output_folder = original_output
            
            results[contrast] = {
                'slide_files': slide_files,
                'output_folder': contrast_output_folder,
                'slide_count': len(slide_files),
                'pptx_file': pptx_file,
                'summary_slide': summary_slide,
                'animated_gif': animated_gif
            }
        
        print(f"\n=== ALL FLIPBOOK SLIDES GENERATED SUCCESSFULLY ===")
        print(f"Generated slides for {len(results)} contrast types: {list(results.keys())}")
        print(f"Slides saved to: {os.path.abspath(self.output_folder)}")
        
        # Print absolute paths for verification
        pptx_count = 0
        summary_count = 0
        gif_count = 0
        for contrast, result in results.items():
            if result['slide_files']:
                extras = []
                if result.get('summary_slide'):
                    extras.append("summary slide")
                    summary_count += 1
                if result.get('animated_gif'):
                    extras.append("animated GIF")
                    gif_count += 1
                if result.get('pptx_file'):
                    extras.append("PowerPoint")
                    pptx_count += 1
                
                extras_text = f" + {', '.join(extras)}" if extras else ""
                print(f"  {contrast}: {len(result['slide_files'])} slides{extras_text} in {os.path.abspath(result['output_folder'])}")
        
        # Print additional feature counts
        features = []
        if summary_count > 0:
            features.append(f"{summary_count} tumor summary slide{'s' if summary_count > 1 else ''}")
        if gif_count > 0:
            features.append(f"{gif_count} animated GIF{'s' if gif_count > 1 else ''}")
        if pptx_count > 0:
            features.append(f"{pptx_count} PowerPoint flipbook{'s' if pptx_count > 1 else ''}")
            
        if features:
            print(f"\n🎉 Enhanced features generated: {', '.join(features)}!")
        
        return results
    
    def create_powerpoint_flipbook(self, slide_files, contrast_type, title="Brain Tumor Flipbook"):
        """Create a PowerPoint presentation from slide files with black background and centered images"""
        if not PPTX_AVAILABLE:
            print("Error: python-pptx not installed. Cannot create PowerPoint file.")
            print("Install with: pip install python-pptx")
            return None
            
        if not slide_files:
            print("No slide files provided for PowerPoint creation")
            return None
        
        from pptx.dml.color import RGBColor
        from PIL import Image as PILImage
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions for widescreen (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Add title slide with black background
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Set black background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black
        
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        seg_info = " with Tumor Overlay" if self.segmentation_folder else ""
        subtitle_placeholder.text = f"{contrast_type} Contrast{seg_info}\nLongitudinal Brain Tumor Analysis\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Make title text white
        title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add slides for each timepoint
        blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
        
        for i, slide_file in enumerate(slide_files):
            if not os.path.exists(slide_file):
                print(f"Warning: Slide file not found: {slide_file}")
                continue
                
            # Extract timepoint from filename
            timepoint_match = re.search(r'(\d+)\.png$', os.path.basename(slide_file))
            timepoint = timepoint_match.group(1) if timepoint_match else str(i+1)
            
            # Create new slide with black background
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Set black background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black
            
            # Add image centered without distortion
            try:
                # Get original image dimensions
                with PILImage.open(slide_file) as img:
                    img_width_px, img_height_px = img.size
                    aspect_ratio = img_width_px / img_height_px
                
                # Calculate slide dimensions
                slide_width = Inches(13.33)
                slide_height = Inches(7.5)
                slide_aspect_ratio = 13.33 / 7.5
                
                # Calculate image size to fit without distortion
                if aspect_ratio > slide_aspect_ratio:
                    # Image is wider than slide aspect ratio - fit to width
                    img_width = slide_width * 0.9  # Leave 10% margin
                    img_height = img_width / aspect_ratio
                else:
                    # Image is taller than slide aspect ratio - fit to height
                    img_height = slide_height * 0.9  # Leave 10% margin
                    img_width = img_height * aspect_ratio
                
                # Center the image
                img_left = (slide_width - img_width) / 2
                img_top = (slide_height - img_height) / 2
                
                slide.shapes.add_picture(slide_file, img_left, img_top, img_width, img_height)
                print(f"Added slide {i+1}: Timepoint {timepoint} (centered, no distortion)")
                
            except Exception as e:
                print(f"Error adding image {slide_file}: {e}")
                continue
        
        # Save PowerPoint file
        pptx_file = os.path.join(self.output_folder, f'flipbook_{contrast_type}.pptx')
        
        try:
            prs.save(pptx_file)
            print(f"PowerPoint flipbook created: {pptx_file}")
            return pptx_file
        except Exception as e:
            print(f"Error saving PowerPoint file: {e}")
            return None
    
    def generate_flipbook_with_pptx(self, contrast_type, rows=3, cols=5, 
                                   window_level=None, window_width=None,
                                   colormap='gray', title=None, tumor_color='red',
                                   tumor_alpha=0.7, show_contour=True, contour_width=2):
        """
        Generate complete flipbook with PowerPoint output
        """
        if title is None:
            seg_suffix = " with Tumor Overlay" if self.segmentation_folder else ""
            title = f"Brain Tumor Flipbook - {contrast_type}{seg_suffix}"
            
        print("=== Brain Tumor Flipbook Generator with PowerPoint Output ===")
        print(f"Input folder: {self.registered_folder}")
        print(f"Segmentation folder: {self.segmentation_folder}")
        print(f"Output folder: {self.output_folder}")
        print(f"Contrast type: {contrast_type}")
        print(f"Tumor visualization: {'Contour' if show_contour else 'Filled overlay'} in {tumor_color}")
        print(f"PowerPoint output: {'Available' if PPTX_AVAILABLE else 'Not available (install python-pptx)'}")
        
        # Find timepoint folders and contrast types
        self.find_timepoint_folders()
        self.find_contrast_types()
        
        # Create slides
        slide_files = self.create_flipbook_slides(
            contrast_type=contrast_type, rows=rows, cols=cols,
            window_level=window_level, window_width=window_width,
            colormap=colormap, tumor_color=tumor_color,
            tumor_alpha=tumor_alpha, show_contour=show_contour,
            contour_width=contour_width
        )
        
        if not slide_files:
            print("Error: No slides were created. Check your contrast type and files.")
            return None, []
        
        # Create PowerPoint flipbook
        pptx_file = None
        if PPTX_AVAILABLE:
            pptx_file = self.create_powerpoint_flipbook(slide_files, contrast_type, title=title)
        
        print(f"\nFlipbook generated successfully!")
        print(f"- {len(slide_files)} slides created")
        if pptx_file:
            print(f"- PowerPoint flipbook: {pptx_file}")
        print(f"- Individual slides saved in: {self.output_folder}")
        
        return pptx_file, slide_files


# Enhanced convenience functions
def create_brain_flipbook_with_segmentation(registered_folder, segmentation_folder, 
                                          contrast_type, output_folder="flipbooks", 
                                          rows=3, cols=5, window_level=None, window_width=None,
                                          colormap='gray', title=None, tumor_color='red',
                                          tumor_alpha=0.7, show_contour=True, contour_width=2):
    """
    Convenience function to create a brain flipbook with tumor segmentation overlay
    
    Parameters:
    - registered_folder: Path to folder containing timepoint subfolders (5885, 5972, 6070, etc.)
    - segmentation_folder: Path to folder containing segmentation maps (same structure)
    - contrast_type: Type of contrast (T1, T2, T1CE, FL, etc.)
    - output_folder: Output folder for flipbook
    - rows, cols: Mosaic grid size (default 3x5 = 15 slices)
    - window_level, window_width: Manual windowing (None for auto)
    - colormap: Matplotlib colormap for brain ('gray', 'bone', etc.)
    - title: Flipbook title
    - tumor_color: Color for tumor ('red', 'yellow', 'cyan', 'green', 'magenta')
    - tumor_alpha: Transparency of tumor overlay (0-1)
    - show_contour: If True, show contour; if False, show filled overlay
    - contour_width: Width of contour lines
    """
    generator = BrainFlipbookGenerator(registered_folder, segmentation_folder, output_folder)
    return generator.generate_flipbook(
        contrast_type=contrast_type, rows=rows, cols=cols,
        window_level=window_level, window_width=window_width,
        colormap=colormap, title=title, tumor_color=tumor_color,
        tumor_alpha=tumor_alpha, show_contour=show_contour, contour_width=contour_width
    )

def create_all_contrast_flipbooks_with_segmentation(registered_folder, segmentation_folder,
                                                  output_folder="all_flipbooks", rows=3, cols=5, 
                                                  window_level=None, window_width=None,
                                                  colormap='gray', tumor_color='red',
                                                  tumor_alpha=0.7, show_contour=True, 
                                                  contour_width=2):
    """
    Create flipbooks for ALL available contrast types with tumor segmentation overlay
    
    Parameters:
    - registered_folder: Path to folder containing timepoint subfolders
    - segmentation_folder: Path to folder containing segmentation maps
    - output_folder: Output folder for all flipbooks (will create subfolders for each contrast)
    - rows, cols: Mosaic grid size
    - window_level, window_width: Manual windowing (None for auto)
    - colormap: Matplotlib colormap for brain
    - tumor_color: Color for tumor overlay
    - tumor_alpha: Transparency of tumor overlay (0-1)
    - show_contour: If True, show contour; if False, show filled overlay
    - contour_width: Width of contour lines
    """
    generator = BrainFlipbookGenerator(registered_folder, segmentation_folder, output_folder)
    return generator.generate_all_contrast_flipbooks(
        rows=rows, cols=cols, window_level=window_level, 
        window_width=window_width, colormap=colormap,
        tumor_color=tumor_color, tumor_alpha=tumor_alpha, 
        show_contour=show_contour, contour_width=contour_width
    )

# Example usage:
if __name__ == "__main__":
    # Example usage with segmentation overlay for ALL contrasts
    registered_folder = r"C:\Users\smita\registered"  # Your registered images folder
    segmentation_folder = r"C:\Users\smita\segmentations"  # Your segmentation folder
    
    # Create flipbooks for ALL available contrasts (T1, T2, T1CE, FL)
    print("=== Creating flipbooks for ALL contrast types ===")
    results = create_all_contrast_flipbooks_with_segmentation(
        registered_folder=registered_folder,
        segmentation_folder=segmentation_folder,
        output_folder="all_contrast_flipbooks",
        tumor_color='red',          # Red tumor overlay
        show_contour=True,          # Show contour instead of filled
        contour_width=2,            # Contour line width
        tumor_alpha=0.8             # Opacity of overlay
    )
    
    print(f"\n*** SUCCESS! Generated flipbooks for: {list(results.keys())}")
    print(f"*** All flipbooks saved in separate folders")
    
    # Alternative: Create individual flipbook for specific contrast
    # slides = create_brain_flipbook_with_segmentation(
    #     registered_folder=registered_folder,
    #     segmentation_folder=segmentation_folder,
    #     contrast_type="T1CE",  # Just T1CE
    #     output_folder="T1CE_only",
    #     tumor_color='yellow',
    #     show_contour=False,     # Filled overlay instead
    #     tumor_alpha=0.5
    # )